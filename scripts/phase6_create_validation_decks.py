from pathlib import Path
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

OUT = Path(__file__).resolve().parents[1] / ".phase6-runtime" / "decks"
OUT.mkdir(parents=True, exist_ok=True)

W, H = Inches(13.333), Inches(7.5)
BG = RGBColor(10, 18, 32)
INK = RGBColor(238, 243, 248)
MUTED = RGBColor(160, 177, 194)
ACCENTS = [RGBColor(45, 212, 191), RGBColor(251, 191, 36), RGBColor(96, 165, 250)]

DECKS = [
    ("AI 课程设计：从知识传递到学习任务", "teacher_ai_lesson_design", [
        ("问题", "教师知道 AI 很重要，但课堂仍停留在工具演示。"),
        ("方法", "把课程拆成真实任务、证据产物和反馈循环。"),
        ("结果", "学生不只会使用工具，还能解释、协作和交付。"),
    ]),
    ("产教融合：把合作机会变成可执行项目", "industry_education_fusion", [
        ("变化", "合作不再只是签约，而要形成共同目标和责任分工。"),
        ("机制", "用项目任务连接学校、企业、教师和学生。"),
        ("落地", "每一阶段都有产物、负责人、证据和复盘。"),
    ]),
    ("教师评价：让 AI 进入课堂后的质量可见", "teacher_ai_assessment", [
        ("难点", "看见了 AI 使用，不等于看见了学习发生。"),
        ("框架", "同时观察过程、作品、解释能力和迁移能力。"),
        ("闭环", "评价结果反过来调整任务设计，而不是只给分数。"),
    ]),
]

def add_text(slide, text, x, y, w, h, size, color=INK, bold=False, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame; tf.clear(); tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = text; r.font.name = "Aptos"; r.font.size = Pt(size); r.font.bold = bold; r.font.color.rgb = color
    return box

def rect(slide, x, y, w, h, fill, radius=False):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid(); shape.fill.fore_color.rgb = fill; shape.line.fill.background()
    return shape

def base(prs, title, kicker, number):
    slide = prs.slides.add_slide(prs.slide_layouts[6]); rect(slide, 0, 0, 13.333, 7.5, BG)
    rect(slide, 0.65, 0.62, 0.08, 0.62, ACCENTS[number % 3])
    add_text(slide, kicker.upper(), 0.95, 0.62, 5.8, 0.28, 11, ACCENTS[number % 3], True)
    add_text(slide, title, 0.95, 1.02, 11.2, 0.7, 28, INK, True)
    add_text(slide, f"TIANSHU / PHASE 6 VALIDATION  ·  {number:02d}", 0.95, 7.08, 5.0, 0.2, 9, MUTED)
    add_text(slide, "真实课题验证样本 · 非 Teacher PPT 069 主仓", 8.2, 7.08, 4.2, 0.2, 9, MUTED, False, PP_ALIGN.RIGHT)
    return slide

def make_deck(title, slug, cards):
    prs = Presentation(); prs.slide_width = W; prs.slide_height = H
    slide = base(prs, title, "三课题产品化验证", 1)
    add_text(slide, "同一总控，三个不同问题，三套不同表达。", 0.95, 2.05, 7.8, 0.45, 19, MUTED)
    for i, (head, body) in enumerate(cards):
        x = 0.95 + i * 4.05; rect(slide, x, 3.1, 3.45, 2.2, RGBColor(20, 34, 53), True)
        rect(slide, x, 3.1, 3.45, 0.08, ACCENTS[i])
        add_text(slide, head, x + 0.28, 3.42, 2.8, 0.35, 16, ACCENTS[i], True)
        add_text(slide, body, x + 0.28, 4.0, 2.8, 0.95, 17, INK)
    slide = base(prs, "一页看懂这个课题的产品逻辑", "结构", 2)
    add_text(slide, "输入", 1.1, 2.35, 1.5, 0.3, 15, ACCENTS[0], True)
    add_text(slide, "方法", 5.35, 2.35, 1.5, 0.3, 15, ACCENTS[1], True)
    add_text(slide, "可见结果", 9.55, 2.35, 2, 0.3, 15, ACCENTS[2], True)
    for x in (2.55, 6.8): rect(slide, x, 2.46, 2.0, 0.05, MUTED)
    for x, t, c in [(1.1, "真实需求\n与场景", ACCENTS[0]), (5.35, "任务拆解\n与协作", ACCENTS[1]), (9.55, "作品、证据\n与反馈", ACCENTS[2])]:
        rect(slide, x, 3.0, 2.55, 1.35, RGBColor(20, 34, 53), True); add_text(slide, t, x + .2, 3.35, 2.15, .65, 18, c, True)
    slide = base(prs, "验收：不是好看，而是能被使用", "验收", 3)
    checks = ["内容与课题一致", "三页逻辑不重复", "视觉层级清楚", "正文可读，无溢出", "输出可打开，可追溯"]
    for i, c in enumerate(checks):
        y = 2.05 + i * .72; rect(slide, 1.0, y, .28, .28, ACCENTS[i % 3], True); add_text(slide, c, 1.55, y - .02, 8.5, .35, 18, INK)
    prs.save(OUT / f"{slug}.pptx")
    Presentation(OUT / f"{slug}.pptx")

for title, slug, cards in DECKS: make_deck(title, slug, cards)
print(f"created {len(DECKS)} decks in {OUT}")
